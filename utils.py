import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

def search_unsplash_image(query, access_key):
    url = f"https://api.unsplash.com/photos/random?query={query}&client_id={access_key}"
    res = requests.get(url)
    if res.status_code == 200:
        data = res.json()
        image_url = data.get('urls', {}).get('regular')
        if image_url:
            img_data = requests.get(image_url).content
            return BytesIO(img_data)
    return None

def create_styled_pptx(slides):
    prs = Presentation()
    for slide in slides:
        s = prs.slides.add_slide(prs.slide_layouts[5])
        title = s.shapes.title
        title.text = slide["title"]
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True

        textbox = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(3))
        tf = textbox.text_frame
        tf.text = slide["content"]
        tf.paragraphs[0].font.size = Pt(20)

        if "image" in slide and slide["image"]:
            try:
                s.shapes.add_picture(slide["image"], Inches(1), Inches(3.5), width=Inches(6))
            except:
                pass

    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

def convert_pptx_to_pdf(pptx_io, output_path="output.pdf"):
    import comtypes.client
    from tempfile import NamedTemporaryFile

    with NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        tmp.write(pptx_io.read())
        tmp_path = tmp.name

    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    presentation = powerpoint.Presentations.Open(tmp_path)
    presentation.SaveAs(output_path, FileFormat=32)
    presentation.Close()
    powerpoint.Quit()
    return output_path
