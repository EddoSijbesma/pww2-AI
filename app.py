import streamlit as st
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from io import BytesIO
import datetime
import os
import requests
from dotenv import load_dotenv

# === ENV ===
load_dotenv()
AZURE_API_KEY = os.getenv("AZURE_OPENAI_KEY")
AZURE_API_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT", "").rstrip("/")
AZURE_MODEL = os.getenv("AZURE_OPENAI_MODEL", "openai/gpt-4.1")

# === AI ASSISTENT VIA AZURE OPENAI ===
def vraag_ai_azure(vraag: str) -> str:
    try:
        response = requests.post(
            f"{AZURE_API_ENDPOINT}/chat/completions",
            headers={
                "Content-Type": "application/json",
                "Authorization": f"Bearer {AZURE_API_KEY}"
            },
            json={
                "messages": [
                    {"role": "system", "content": ""},
                    {"role": "user", "content": vraag}
                ],
                "model": AZURE_MODEL,
                "temperature": 0.7,
                "top_p": 1
            },
            timeout=15
        )
        response.raise_for_status()
        return response.json()["choices"][0]["message"]["content"]
    except Exception as e:
        return f"⚠️ Fout bij AI-aanvraag: {e}"

# === TITEL EN AI-SIDEBAR ===
st.title("Stappenplan Maker Gemaakt door Eddo.S")

st.sidebar.header("🤖 AI Assistent")
user_prompt = st.sidebar.text_area("Stel je vraag aan de AI:")

if user_prompt.strip() != "":
    with st.sidebar.spinner("AI denkt na..."):
        antwoord = vraag_ai_azure(user_prompt)
        st.sidebar.markdown("### Antwoord:")
        st.sidebar.write(antwoord)

# === KLEURTHEMA ===
kleur = st.selectbox("Kies een kleurthema", ["Blauw", "Groen", "Rood", "Grijs"])
kleurmap = {
    "Blauw": RGBColor(0, 112, 192),
    "Groen": RGBColor(0, 176, 80),
    "Rood": RGBColor(192, 0, 0),
    "Grijs": RGBColor(89, 89, 89)
}
geselecteerde_kleur = kleurmap.get(kleur, RGBColor(0, 112, 192))

# === GEGEVENS ===
st.header("📋 Gegevens")
naam = st.text_input("Naam student")
studentnummer = st.text_input("Studentnummer")
projectnaam = st.text_input("Naam project")
locatie = st.text_input("Locatie project")
leerbedrijf = st.text_input("Leerbedrijf")
leermeester = st.text_input("Leermeester")
inleverdatum = st.date_input("Inleverdatum", datetime.date.today())

st.file_uploader("Upload hier foto's van jezelf tijdens het werk", accept_multiple_files=True)

# === PRAKTIJKOPDRACHT ===
st.header("🛠️ Over de praktijkopdracht")
opdracht = st.text_area("Welke praktijkopdracht heb je gemaakt?")
wat_gemaakt = st.text_area("Wat heb je gemaakt?")
waarom_gemaakt = st.text_area("Waarom heb je deze praktijkopdracht gekozen?")
type_werk = st.selectbox("Wat voor type werk was het?", ["Nieuwbouw", "Aanbouw", "Renovatie", "Onderhoud", "Anders"])
werksituatie = st.text_area("Hoe was de werksituatie? (bijv. samenwerking, tijdsdruk, weer)")
ploeggrootte = st.text_input("Hoe groot was je ploeg?")
st.file_uploader("Upload hier foto’s van het eindresultaat", accept_multiple_files=True)

# === RISICO ===
st.header("⚠️ Risico’s en maatregelen")
risicos = st.text_area("Beschrijf de risico’s bij deze praktijkopdracht")
maatregelen = st.text_area("Welke maatregelen heb je getroffen?")

# === WERKTEKENING ===
st.header("📐 Werktekening")
st.file_uploader("Upload hier je werktekening", type=["jpg", "png", "pdf"])

# === MATERIAAL EN GEREEDSCHAP ===
st.header("🧰 Materiaal en gereedschap")
materialen = st.text_area("Materiaalstaat")
gereedschap = st.text_area("Gereedschapslijst")
werkuur = st.text_area("Werkschema en urenverantwoording")

# === STAPPENPLAN ===
st.header("🪜 Stappenplan")
for i in range(1, 11):
    with st.expander(f"Stap {i}"):
        st.text_input(f"Stap {i} – Titel", key=f"stap{i}_titel")
        st.text_area(f"Wat heb je gedaan?", key=f"stap{i}_wat")
        st.text_area(f"Waarom heb je het zo gedaan?", key=f"stap{i}_waarom")
        st.text_area(f"Wat was een leerpunt?", key=f"stap{i}_leer")
        st.text_area(f"Instructies voor je collega", key=f"stap{i}_instructie")
        st.text_area(f"Let op!", key=f"stap{i}_letop")
        st.file_uploader("Voeg hier foto's toe", accept_multiple_files=True, key=f"stap{i}_foto")

# === REFLECTIE PERSOONLIJK ===
st.header("🔍 Reflectie: Persoonlijk")
st.text_area("Hoeveel hulp had je nodig en wat kon je zelfstandig?", key="reflectie_hulp")
st.text_area("Wanneer stuurde een collega je bij?", key="reflectie_bijsturing")
st.text_area("Welke tips heb je gekregen?", key="reflectie_tips")
st.text_area("Wat waren je leerpunten?", key="reflectie_leerpunt")
st.text_area("Wat waren je sterke punten?", key="reflectie_sterk")

# === REFLECTIE SAMENWERKEN ===
st.header("👥 Reflectie: Samenwerken")
st.text_area("Wat werd er van je verwacht?", key="samen_verwacht")
st.text_area("Wat heb je zelfstandig gedaan?", key="samen_zelfstandig")
st.text_area("Wat zou je de volgende keer anders doen?", key="samen_anders")
st.text_area("Wat wil je nog leren?", key="samen_leren")
st.text_area("Tips van je collega?", key="samen_tips")

# === POWERPOINT GENERATIE ===
def vervang_tekst(tekst, vervangingen):
    for sleutel, waarde in vervangingen.items():
        tekst = tekst.replace(f"{{{{{sleutel}}}}}", str(waarde))
    return tekst

def genereer_powerpoint(vervangingen, template_path="tamplatepraktijkopdracht2.pptx"):
    prs = Presentation(template_path)
    for slide in prs.slides:
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            titel = slide.shapes.title
            titel.text = vervang_tekst(titel.text, vervangingen)
            for para in titel.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = geselecteerde_kleur
                    run.font.size = Pt(20)
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                shape.text = vervang_tekst(shape.text, vervangingen)
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = geselecteerde_kleur
                        run.font.size = Pt(12)
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# === DOWNLOAD ===
st.header("📤 Afronding en Download")

vervangingen = {
    "naam": naam,
    "studentnummer": studentnummer,
    "projectnaam": projectnaam,
    "locatie": locatie,
    "leerbedrijf": leerbedrijf,
    "leermeester": leermeester,
    "inleverdatum": inleverdatum.strftime('%d-%m-%Y'),
    "opdracht": opdracht,
    "wat_gemaakt": wat_gemaakt,
    "waarom_gemaakt": waarom_gemaakt,
    "type_werk": type_werk,
    "werksituatie": werksituatie,
    "ploeggrootte": ploeggrootte,
    "risicos": risicos,
    "maatregelen": maatregelen,
    "materialen": materialen,
    "gereedschap": gereedschap,
    "werkuur": werkuur,
    "reflectie_hulp": st.session_state.get("reflectie_hulp", ""),
    "reflectie_bijsturing": st.session_state.get("reflectie_bijsturing", ""),
    "reflectie_tips": st.session_state.get("reflectie_tips", ""),
    "reflectie_leerpunt": st.session_state.get("reflectie_leerpunt", ""),
    "reflectie_sterk": st.session_state.get("reflectie_sterk", ""),
    "samen_verwacht": st.session_state.get("samen_verwacht", ""),
    "samen_zelfstandig": st.session_state.get("samen_zelfstandig", ""),
    "samen_anders": st.session_state.get("samen_anders", ""),
    "samen_leren": st.session_state.get("samen_leren", ""),
    "samen_tips": st.session_state.get("samen_tips", ""),
    "kleur": kleur
}

for i in range(1, 11):
    vervangingen[f"stap{i}_titel"] = st.session_state.get(f"stap{i}_titel", "")
    vervangingen[f"stap{i}_wat"] = st.session_state.get(f"stap{i}_wat", "")
    vervangingen[f"stap{i}_waarom"] = st.session_state.get(f"stap{i}_waarom", "")
    vervangingen[f"stap{i}_leer"] = st.session_state.get(f"stap{i}_leer", "")
    vervangingen[f"stap{i}_instructie"] = st.session_state.get(f"stap{i}_instructie", "")
    vervangingen[f"stap{i}_letop"] = st.session_state.get(f"stap{i}_letop", "")

if st.button("📥 Genereer & Download PowerPoint (.pptx)"):
    try:
        bestand = genereer_powerpoint(vervangingen)
        st.download_button(
            label="⬇️ Download PowerPoint",
            data=bestand,
            file_name="praktijkopdracht.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    except Exception as e:
        st.error(f"Er is een fout opgetreden: {e}")
